"""生成された pptx が仕様を満たすことの検証（生成物を開き直して確認）。"""
import sys
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts" / "seminar_p77"))

import content_p77 as C
from build_p77_slides import build


@pytest.fixture(scope="module")
def deck_path(tmp_path_factory):
    out = tmp_path_factory.mktemp("deck") / "p77_test.pptx"
    build(str(out))
    return out


@pytest.fixture(scope="module")
def prs(deck_path):
    return Presentation(str(deck_path))


def iter_runs(slide):
    """スライド上の全 run（テキストボックス＋表セル）を列挙。"""
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                yield from para.runs
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        yield from para.runs


def slide_text(slide):
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            texts.append(sh.text_frame.text)
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    texts.append(cell.text_frame.text)
    return "\n".join(texts)


def test_slide_count(prs):
    assert len(prs.slides._sldIdLst) == 8


def test_canvas_is_16_9(prs):
    assert prs.slide_width == 12192000 and prs.slide_height == 6858000


def test_all_runs_at_least_20pt_and_explicit(prs):
    for i, slide in enumerate(prs.slides):
        for run in iter_runs(slide):
            if not run.text.strip():
                continue
            assert run.font.size is not None, \
                f"slide {i+1}: サイズ未指定の run 「{run.text[:20]}」"
            assert run.font.size >= Pt(20), \
                f"slide {i+1}: {run.font.size.pt}pt < 20pt 「{run.text[:20]}」"


def test_fonts_are_ibm_plex(prs):
    for i, slide in enumerate(prs.slides):
        for run in iter_runs(slide):
            if not run.text.strip():
                continue
            assert run.font.name and run.font.name.startswith("IBM Plex"), \
                f"slide {i+1}: フォント {run.font.name} 「{run.text[:20]}」"


def test_slide_texts_match_content_module(prs):
    for i, slide in enumerate(prs.slides):
        built = slide_text(slide)
        assert C.SLIDES[i]["headline"] in built, f"slide {i+1}: 見出しが描画されていない"


def test_key_strings_rendered(prs):
    slides = list(prs.slides)
    assert "Model Context Protocol" in slide_text(slides[2])
    assert "60" in slide_text(slides[4])
    assert "backup ファイル" in slide_text(slides[5])
    assert "抄録のみ" in slide_text(slides[7])


def test_notes_attached_verbatim(prs):
    for i, slide in enumerate(prs.slides):
        notes = slide.notes_slide.notes_text_frame.text
        assert notes == C.SLIDES[i]["notes"], f"slide {i+1}: 台本がノート欄と不一致"


def test_package_part_names_unique_and_only_eight_slides(deck_path):
    names = zipfile.ZipFile(deck_path).namelist()
    assert len(names) == len(set(names)), "zip 内に同名パーツが重複している"
    slide_parts = [n for n in names
                   if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    assert len(slide_parts) == 8, f"スライドパーツが {len(slide_parts)} 個（8 個のはず）"


from pptx.enum.shapes import MSO_SHAPE_TYPE


def count_autoshapes(slide):
    return sum(1 for sh in slide.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)


def test_figures_present(prs):
    slides = list(prs.slides)
    # S3: 箱3+矢印2 ≥ 5 / S4: 箱5+矢印4+戻り矢印1 ≥ 9（他要素は加算されるだけ）
    # S6: 箱・矢印・結果箱×2レーン+下部バー ≥ 7
    assert count_autoshapes(slides[2]) >= 5, "S3 の MCP 概念図がない"
    assert count_autoshapes(slides[3]) >= 9, "S4 の構成図がない"
    assert count_autoshapes(slides[5]) >= 7, "S6 の回避ルート図がない"


def test_figure_labels_rendered(prs):
    slides = list(prs.slides)
    assert "共通の差し込み口" in slide_text(slides[2])
    assert "Aspen Plus" in slide_text(slides[3])
    assert "窓口がない" in slide_text(slides[5])
