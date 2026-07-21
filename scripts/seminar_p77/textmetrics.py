# -*- coding: utf-8 -*-
"""実フォント計測にもとづくテキスト折り返し・高さ計算とレイアウト監査。

スライドで実際に使うフォント（欧文 IBM Plex Sans / 和文 Hiragino Kaku Gothic
ProN）のグリフ幅を matplotlib 経由で測り、PowerPoint と同じ word-wrap 規則
（欧単語は分割しない・和文は文字単位で折り返し可）で行数を数える。
builder は本モジュールの値で要素を「積み上げ」配置し、テストは
audit_presentation() で生成物のはみ出し・重なりゼロを検証する。
"""
import os
from functools import lru_cache

# PowerPoint のテキストボックス既定インセット（左右 0.1in ずつ）
INSET = 0.1
# 計測誤差・カーニング差の安全率
SAFETY = 1.05

_LATIN_REG = os.path.expanduser("~/Library/Fonts/IBMPlexSans-Regular.ttf")
_LATIN_BOLD = os.path.expanduser("~/Library/Fonts/IBMPlexSans-SemiBold.ttf")
_JP_REG = "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc"
_JP_BOLD = "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc"

try:
    from matplotlib.font_manager import FontProperties
    from matplotlib.textpath import TextToPath

    _T2P = TextToPath()
    _HAVE_MPL = all(os.path.exists(p) for p in (_LATIN_REG, _LATIN_BOLD, _JP_REG))
except Exception:  # matplotlib なし等
    _HAVE_MPL = False

if not os.path.exists(_JP_BOLD):
    _JP_BOLD = _JP_REG  # 和文はウェイトで送り幅が変わらないため W3 で代用可


def _is_jp(ch: str) -> bool:
    """全角として扱う文字か（CJK・かな・全角記号）。"""
    o = ord(ch)
    return (
        0x3000 <= o <= 0x30FF      # CJK 記号・かな
        or 0x4E00 <= o <= 0x9FFF   # 漢字
        or 0xFF00 <= o <= 0xFFEF   # 全角英数・記号
        or ch in "―…〜※→↔①②③"
    )


@lru_cache(maxsize=8192)
def _measure_token(token: str, size_pt: float, bold: bool, jp: bool) -> float:
    """token の幅 [in]。同一スクリプトのみからなる token を渡すこと。"""
    if _HAVE_MPL:
        path = (_JP_BOLD if bold else _JP_REG) if jp else (
            _LATIN_BOLD if bold else _LATIN_REG)
        w, _h, _d = _T2P.get_text_width_height_descent(
            token, FontProperties(fname=path, size=size_pt), False)
        return w / 72.0 * SAFETY
    # フォールバック（フォント不在環境）: 全角=1.0em, 半角=0.55em
    em = size_pt / 72.0
    return sum((1.0 if _is_jp(ch) else 0.55) for ch in token) * em * SAFETY


def _tokenize(text: str, bold: bool):
    """折り返し単位の列 [(width_fn 引数タプル…)] を返す。

    欧文は空白区切りの単語ごと（分割不可）、和文は 1 文字ごと（分割可）。
    戻り値: list of (token, jp, breaking) — breaking は直前で改行してよいか。
    """
    tokens = []
    buf = ""
    for ch in text:
        if _is_jp(ch):
            if buf:
                tokens.append((buf, False))
                buf = ""
            tokens.append((ch, True))
        elif ch == " ":
            if buf:
                tokens.append((buf, False))
                buf = ""
            tokens.append((" ", False))
        else:
            buf += ch
    if buf:
        tokens.append((buf, False))
    return [(t, jp, bold) for t, jp in tokens]


def wrap_lines(runs, usable_w_in: float, size_pt: float) -> int:
    """runs=[(text, bold), ...] を usable 幅で折り返したときの行数。"""
    tokens = []
    for text, bold in runs:
        for para_piece in [text]:  # \n は呼び出し側で分割済みの想定
            tokens += _tokenize(para_piece, bold)
    lines, cur = 1, 0.0
    for tok, jp, bold in tokens:
        w = _measure_token(tok, size_pt, bold, jp)
        if cur + w <= usable_w_in or cur == 0.0:
            cur += w
        else:
            lines += 1
            cur = 0.0 if tok == " " else w  # 行頭の空白は消える
    return lines


def para_height_in(runs, box_w_in: float, size_pt: float,
                   spacing: float, space_after_pt: float = 0.0) -> float:
    """1 段落の描画高さ [in]（折り返し行数 × 行送り + 段落後間隔）。"""
    n = wrap_lines(runs, box_w_in - 2 * INSET, size_pt)
    return n * size_pt * spacing / 72.0 + space_after_pt / 72.0


def block_height_in(paragraphs, box_w_in: float) -> float:
    """複数段落ブロックの高さ。paragraphs=[(runs, size, spacing, space_after_pt)]"""
    return sum(para_height_in(r, box_w_in, s, sp, sa)
               for r, s, sp, sa in paragraphs)


# ---------------------------------------------------------------- 監査

def _shape_paragraph_specs(sh):
    """pptx shape からテキスト段落仕様を復元して [(runs,size,spacing,sa)] を返す。"""
    specs = []
    for para in sh.text_frame.paragraphs:
        runs = [(r.text, bool(r.font.bold)) for r in para.runs if r.text]
        if not runs:
            continue
        size = max((r.font.size.pt for r in para.runs if r.font.size), default=20)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        sa = para.space_after.pt if para.space_after is not None else 0.0
        specs.append((runs, size, spacing, sa))
    return specs


def _emu_in(v) -> float:
    return v / 914400.0


def audit_presentation(prs, canvas_h: float = 7.5, canvas_w: float = 13.3334):
    """全スライドを検査し、違反メッセージのリストを返す（空なら合格）。

    検査項目:
      (a) テキスト実高がキャンバス下端 7.45in を超えない
      (b) 図形（中央アンカー）のテキスト実高が図形高さに収まる
      (c) テキスト矩形どうしの重なりなし（空テキストのパネル・矢印は除外）
    """
    from pptx.enum.text import MSO_ANCHOR

    problems = []
    for si, slide in enumerate(prs.slides, start=1):
        rects = []  # (label, x0, y0, x1, y1)
        for sh in slide.shapes:
            x, y = _emu_in(sh.left), _emu_in(sh.top)
            w, h = _emu_in(sh.width), _emu_in(sh.height)
            if getattr(sh, "has_table", False) and sh.has_table:
                th = sum(_emu_in(r.height) for r in sh.table.rows)
                rects.append((f"table@{y:.2f}", x, y, x + w, y + max(th, h)))
                continue
            if not sh.has_text_frame:
                continue
            specs = _shape_paragraph_specs(sh)
            if not specs:
                continue  # 空テキスト（背景パネル・矢印）は対象外
            text_h = block_height_in(specs, w) + 0.04
            label = f"「{specs[0][0][0][0][:12]}…」@{y:.2f}"
            anchored_mid = (sh.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE)
            if anchored_mid:
                if text_h > h + 0.03:
                    problems.append(
                        f"S{si} (b) 図形内テキストあふれ {label}: "
                        f"必要 {text_h:.2f}in > 箱 {h:.2f}in")
                ty0 = y + max(0.0, (h - text_h) / 2)
            else:
                ty0 = y
            ty1 = ty0 + text_h
            if ty1 > canvas_h - 0.05:
                problems.append(
                    f"S{si} (a) 下端はみ出し {label}: 底 {ty1:.2f}in")
            if _emu_in(sh.left) + _emu_in(sh.width) > canvas_w + 0.01:
                problems.append(f"S{si} 右端はみ出し {label}")
            rects.append((label, x, ty0, x + w, ty1))
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                la, ax0, ay0, ax1, ay1 = rects[i]
                lb, bx0, by0, bx1, by1 = rects[j]
                ox = min(ax1, bx1) - max(ax0, bx0)
                oy = min(ay1, by1) - max(ay0, by0)
                if ox > 0.05 and oy > 0.02:
                    problems.append(
                        f"S{si} (c) 重なり {la} × {lb}: {ox:.2f}×{oy:.2f}in")
    return problems
