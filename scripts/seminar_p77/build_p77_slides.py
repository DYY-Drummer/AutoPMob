# -*- coding: utf-8 -*-
"""MCP 活用研究 紹介スライド（v2・10 枚）のビルダー。

テンプレ（PSE Asia 2026 発表スライド）を複製し、既存スライドを全削除して
content_p77.SLIDES の 10 枚を描画する。
各要素は textmetrics の実フォント計測にもとづく「積み上げ」配置
（折り返し行数を実測して次要素の y を決める）で、はみ出し・重なりを防ぐ。
使い方: python3 scripts/seminar_p77/build_p77_slides.py [出力パス]
"""
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import content_p77 as C
import textmetrics as TM

TEMPLATE = "/Users/kazuhiromiyamura/Desktop/学会/PSE ASIAN2026/PSEAsia2026_slide_v2.pptx"
DEFAULT_OUT = "/Users/kazuhiromiyamura/Desktop/ゼミ/論文紹介_LLM_MCP_シミュレータ操作_かずひろ.pptx"

NAVY = RGBColor(0x16, 0x30, 0x4F)
TEXT = RGBColor(0x22, 0x22, 0x22)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
GREEN = RGBColor(0x2E, 0x88, 0x51)
GRAY = RGBColor(0x8C, 0x8C, 0x8C)
LIGHT = RGBColor(0xEA, 0xF3, 0xFB)
LIGHT2 = RGBColor(0xDC, 0xEA, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BODY_FONT = "IBM Plex Sans"
HEAD_FONT = "IBM Plex Sans SemiBold"

CANVAS_W = 13.3334
MARGIN_X = 0.55
FULL_W = CANVAS_W - 2 * MARGIN_X  # 12.23


def _style(run, size, color, bold, font):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font  # 欧文フォント（<a:latin>）
    # IBM Plex に和文グリフは無いため、和文フォント（<a:ea>）を明示指定する
    rPr = run.font._rPr
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "Hiragino Kaku Gothic ProN")


def add_text(slide, x, y, w, lines, size=20, color=TEXT, bold=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, spacing=1.15):
    """テキストボックスを置き、実測高さから下端 y を返す。lines は str か list[str]。"""
    if isinstance(lines, str):
        lines = [lines]
    text_h = sum(TM.para_height_in([(ln, bold)], w, size, spacing)
                 for ln in lines) + 0.04
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(text_h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = line
        _style(run, size, color, bold, font)
    return y + text_h


def add_bullets(slide, x, y, w, items, size=20, spacing=1.2):
    """箇条書き。(lead, rest) tuple は lead を太字 NAVY に。下端 y を返す。"""
    sa = 8.0  # space_after [pt]
    text_h = 0.0
    for item in items:
        runs = ([("・", False), (item[0], True), (item[1], False)]
                if isinstance(item, tuple) else [("・" + item, False)])
        text_h += TM.para_height_in(runs, w, size, spacing, sa)
    text_h += 0.04
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(text_h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = spacing
        para.space_after = Pt(sa)
        bullet = para.add_run()
        bullet.text = "・"
        _style(bullet, size, TEXT, False, BODY_FONT)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = para.add_run()
            r1.text = lead
            _style(r1, size, NAVY, True, HEAD_FONT)
            r2 = para.add_run()
            r2.text = rest
            _style(r2, size, TEXT, False, BODY_FONT)
        else:
            r = para.add_run()
            r.text = item
            _style(r, size, TEXT, False, BODY_FONT)
    return y + text_h


def add_box(slide, x, y, w, lines, fill, text_color=WHITE, size=20,
            bold=False, font=BODY_FONT, outline=None, h=None):
    """角丸四角＋中央揃えテキスト。h 省略時はテキスト実測で自動高さ。下端 y を返す。"""
    if isinstance(lines, str):
        lines = lines.split("\n")
    if h is None:
        h = sum(TM.para_height_in([(ln, bold)], w, size, 1.15)
                for ln in lines if ln) + 0.18
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.line_spacing = 1.15
        run = para.add_run()
        run.text = line
        _style(run, size, text_color, bold, font)
    return y + h


def add_arrow(slide, x, y, w, h, fill=GRAY, left=False):
    """矢印（既定は右向き）。"""
    kind = MSO_SHAPE.LEFT_ARROW if left else MSO_SHAPE.RIGHT_ARROW
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_headline(slide, text):
    """見出しを置き、下端 y を返す。"""
    return add_text(slide, MARGIN_X, 0.3, FULL_W, text, size=28, color=NAVY,
                    bold=True, font=HEAD_FONT, spacing=1.1)


# ---- kind 別レンダラ --------------------------------------------------

def render_title(slide, c):
    y = add_text(slide, 0.8, 0.55, 11.7, c["headline"], size=36, color=NAVY,
                 bold=True, font=HEAD_FONT, spacing=1.1)
    y = add_text(slide, 0.8, y + 0.20, 11.7, c["sub"], size=24, color=BLUE,
                 bold=True, font=HEAD_FONT)
    y = add_text(slide, 0.8, y + 0.30, 11.7, c["body"], size=22, spacing=1.25)
    add_text(slide, 0.8, y + 0.25, 11.7, c["explain"], size=20, color=GRAY)
    add_text(slide, 0.8, 6.95, 11.7, c["footer"], size=22, color=TEXT)


def render_table(slide, c):
    hb = add_headline(slide, c["headline"])
    lead_b = add_text(slide, MARGIN_X, hb + 0.12, FULL_W, c["lead"],
                      size=22, color=TEXT, bold=True)
    widths = c.get("col_widths", [4.9, 1.1, 6.25])
    table_y = lead_b + 0.15
    # 各行の高さをセル内折り返しの実測から決める
    row_hs = []
    for row in [c["header"]] + c["rows"]:
        need = 0.0
        for j, val in enumerate(row):
            n = TM.wrap_lines([(val, row is c["header"])], widths[j] - 0.2, 20)
            need = max(need, n * 20 / 72.0)
        row_hs.append(need + 0.14)
    table_shape = slide.shapes.add_table(
        len(row_hs), len(widths), Inches(MARGIN_X), Inches(table_y),
        Inches(FULL_W), Inches(sum(row_hs)))
    table = table_shape.table
    for j, w in enumerate(widths):
        table.columns[j].width = Inches(w)
    for i, rh in enumerate(row_hs):
        table.rows[i].height = Inches(rh)
    for j, head in enumerate(c["header"]):
        _fill_cell(table.cell(0, j), head, bold=True, color=WHITE, fill=NAVY)
    for i, row in enumerate(c["rows"], start=1):
        fill = LIGHT if i % 2 == 1 else WHITE  # 既定の表スタイルを上書きして縞にする
        for j, val in enumerate(row):
            _fill_cell(table.cell(i, j), val, fill=fill,
                       align=PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT)
    note_y = table_y + sum(row_hs) + 0.22
    nb = add_text(slide, MARGIN_X, note_y, FULL_W, c["note"], size=20,
                  color=BLUE, bold=True)
    add_text(slide, MARGIN_X, nb + 0.10, FULL_W, c["source"], size=20, color=GRAY)


def _fill_cell(cell, text, bold=False, color=TEXT, fill=None, align=PP_ALIGN.LEFT):
    cell.margin_top = cell.margin_bottom = Pt(2)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    para.line_spacing = 1.0
    run = para.add_run()
    run.text = text
    _style(run, 20, color, bold, BODY_FONT)


def render_mcp_fig(slide, c):
    hb = add_headline(slide, c["headline"])
    bb = add_bullets(slide, MARGIN_X, hb + 0.25, FULL_W, c["bullets"], size=22)
    fig = c["fig"]
    fy = max(bb + 0.35, 4.15)  # 図ブロック基準 y（tools 箱の上端）
    add_box(slide, 0.7, fy + 0.25, 3.4, fig["llm"], NAVY, WHITE, size=22,
            bold=True, font=HEAD_FONT, h=1.6)
    add_arrow(slide, 4.25, fy + 0.75, 0.75, 0.55)
    add_box(slide, 5.15, fy + 0.40, 2.7, fig["mcp"], BLUE, WHITE, size=20,
            bold=True, h=1.3)
    add_arrow(slide, 8.0, fy + 0.75, 0.75, 0.55)
    add_box(slide, 8.9, fy, 3.7, [fig["tools_title"]] + fig["tools"].split("\n"),
            LIGHT, TEXT, size=20, h=2.1)


def render_arch_fig(slide, c):
    hb = add_headline(slide, c["headline"])
    by = hb + 0.30  # 箱の上端
    # 幅は 20pt 太字（Hiragino W6/IBM Plex SemiBold）の実測折り返し幅から逆算
    xs = [(0.55, 1.95), (2.68, 2.62), (5.48, 2.44), (8.10, 1.95), (10.23, 2.55)]
    fills = [LIGHT2, NAVY, BLUE, LIGHT2, GREEN]
    colors = [TEXT, WHITE, WHITE, TEXT, WHITE]
    for (x, w), fill, col, label in zip(xs, fills, colors, c["boxes"]):
        add_box(slide, x, by, w, label, fill, col, size=20, bold=True, h=1.35)
    for gap_x in (2.50, 5.30, 7.92, 10.05):
        add_arrow(slide, gap_x, by + 0.47, 0.18, 0.42)
    add_arrow(slide, 3.0, by + 1.60, 8.5, 0.5, left=True)
    lb = add_text(slide, 3.0, by + 2.14, 8.5, c["arrow_back"], size=20,
                  color=GRAY, align=PP_ALIGN.CENTER)
    add_bullets(slide, MARGIN_X, lb + 0.18, FULL_W, c["bullets"], size=20)


def render_two_col(slide, c):
    hb = add_headline(slide, c["headline"])
    panel_y = hb + 0.25
    col_w, gap = 5.95, 0.35
    lx, rx = MARGIN_X, MARGIN_X + col_w + gap
    # 中身の高さを先に実測してパネル高さを決める
    def col_h(items):
        h = 0.0
        for item in items:
            runs = ([("・", False), (item[0], True), (item[1], False)]
                    if isinstance(item, tuple) else [("・" + item, False)])
            h += TM.para_height_in(runs, col_w - 0.6, 20, 1.2, 8.0)
        return h
    inner_h = max(col_h(c["left_items"]), col_h(c["right_items"]))
    panel_h = 0.20 + 0.50 + 0.12 + inner_h + 0.22
    add_box(slide, lx, panel_y, col_w, [""], LIGHT2, h=panel_h)
    add_box(slide, rx, panel_y, col_w, [""], LIGHT, h=panel_h)
    for x, title, color, items in (
            (lx, c["left_title"], GRAY, c["left_items"]),
            (rx, c["right_title"], GREEN, c["right_items"])):
        tb = add_text(slide, x + 0.3, panel_y + 0.20, col_w - 0.6, title,
                      size=24, color=color, bold=True, font=HEAD_FONT)
        add_bullets(slide, x + 0.3, tb + 0.12, col_w - 0.6, items, size=20)
    add_box(slide, MARGIN_X, panel_y + panel_h + 0.25, FULL_W, c["bottom"],
            NAVY, WHITE, size=22, bold=True, font=HEAD_FONT)


def render_bullets_box(slide, c):
    hb = add_headline(slide, c["headline"])
    bb = add_bullets(slide, MARGIN_X, hb + 0.30, FULL_W, c["bullets"], size=22,
                     spacing=1.25)
    box_y = max(bb + 0.30, 5.45)
    add_box(slide, MARGIN_X, box_y, FULL_W, c["box"], LIGHT, NAVY, size=24,
            bold=True, font=HEAD_FONT, outline=BLUE)


def render_summary(slide, c):
    hb = add_headline(slide, c["headline"])
    items = [f"{i+1}.  {t}" for i, t in enumerate(c["items"])]
    add_text(slide, 0.8, hb + 0.35, 11.7, items, size=22, spacing=1.45)
    note_h = TM.para_height_in([(c["note"], False)], 11.7, 20, 1.15) + 0.04
    add_text(slide, 0.8, 7.42 - note_h, 11.7, c["note"], size=20, color=GRAY)


RENDERERS = {
    "title": render_title,
    "table": render_table,
    "mcp_fig": render_mcp_fig,
    "arch_fig": render_arch_fig,
    "two_col": render_two_col,
    "bullets_box": render_bullets_box,
    "summary": render_summary,
}


def build(out_path: str) -> None:
    import copy

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE, out)          # マスター・テーマを丸ごと引き継ぐ
    prs = Presentation(str(out))

    # Get reference to first slide's notes structure before deletion
    first_slide = next(iter(prs.slides))
    reference_notes_element = first_slide.notes_slide._element

    sld_list = prs.slides._sldIdLst      # 既存14枚を全削除
    for sld in list(sld_list):
        # rel も切断しないと旧スライドパーツがパッケージ内に残り、
        # 保存時に slide1..N のパート名衝突（Duplicate name 警告）を起こす。
        prs.part.drop_rel(sld.get(qn("r:id")))
        sld_list.remove(sld)

    for c in C.SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # 'DEFAULT'

        # Clone notes structure so notes_placeholder exists
        notes_slide = slide.notes_slide
        reference_cSld = reference_notes_element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
        if reference_cSld is not None:
            new_cSld_elem = copy.deepcopy(reference_cSld)
            existing_cSld = notes_slide._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
            if existing_cSld is not None:
                notes_slide._element.remove(existing_cSld)
            notes_slide._element.insert(0, new_cSld_elem)

        RENDERERS[c["kind"]](slide, c)
        set_notes(slide, c["notes"])
    # テンプレ由来の文書メタデータを本デッキ用に上書き（作成者はそのまま）
    prs.core_properties.title = "論文紹介: LLMにシミュレータを操作させる研究の構造（DTU論文中心）"
    prs.core_properties.subject = "研究室ゼミ（2026-07）"
    prs.save(str(out))


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if a != "--force"]
    force = "--force" in sys.argv[1:]
    target = args[0] if args else DEFAULT_OUT
    # 2026-07-21 以降、既定出力先は本人の手動編集を含む「正本」。
    # 誤って再生成で上書きしないよう、明示の --force がない限り中止する。
    if target == DEFAULT_OUT and Path(target).exists() and not force:
        sys.exit("中止: 出力先は手動編集済みの正本です。上書きするには --force を付けてください。")
    build(target)
    print(f"wrote {target}")
